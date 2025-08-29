# create_master_template.py
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

def create_master_template(filepath="business.pptx"):
    """
    Generates a final, all-in-one 7-slide template that includes every
    special layout type: content, comparison, chart, and table.
    """
    print(f"Creating a final, all-in-one template at: {filepath}")
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.name = 'Title'
    slide.placeholders[1].name = 'Subtitle'

    # Slide 2: Standard Content (Overview)
    slide = prs.slides.add_slide(prs.slide_layouts[8]) # Picture with Caption
    slide.shapes.title.name = 'Title'
    for shape in slide.placeholders:
        if 'Picture' in shape.name: shape.name = 'SlideImage'
        elif 'Text' in shape.name: shape.name = 'Content'

    # Slide 3: Comparison Slide
    slide = prs.slides.add_slide(prs.slide_layouts[3]) # Comparison Layout
    slide.shapes.title.name = 'Title'
    slide.placeholders[1].name = 'ContentLeft'
    slide.placeholders[2].name = 'ContentRight'

    # Slide 4: Chart Slide
    slide = prs.slides.add_slide(prs.slide_layouts[1]) # Title and Content
    slide.shapes.title.name = 'Title'
    slide.placeholders[1].name = 'ChartImage'

    # Slide 5: Table Slide
    slide = prs.slides.add_slide(prs.slide_layouts[5]) # Title Only
    slide.shapes.title.name = 'Title'
    # Add a table shape and name it 'DataTable'
    rows, cols = 5, 4
    table_shape = slide.shapes.add_table(rows, cols, Inches(1), Inches(1.5), Inches(11.33), Inches(4))
    table_shape.name = 'DataTable'

    # Slides 6 & 7: Standard Content
    for _ in range(2):
        slide = prs.slides.add_slide(prs.slide_layouts[8]) # Picture with Caption
        slide.shapes.title.name = 'Title'
        for shape in slide.placeholders:
            if 'Picture' in shape.name: shape.name = 'SlideImage'
            elif 'Text' in shape.name: shape.name = 'Content'

    prs.save(filepath)
    print(f"✅ Master template '{filepath}' created successfully.")

if __name__ == "__main__":
    create_master_template()