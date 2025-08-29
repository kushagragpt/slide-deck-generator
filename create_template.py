# create_template.py
from pptx import Presentation
from pptx.util import Inches


def create_professional_template(filepath="professional_template.pptx"):
    """
    Generates a professional PowerPoint template by adding slides with
    pre-defined layouts and renaming their placeholders for script compatibility.
    """
    print(f"Creating a new professional template at: {filepath}")
    prs = Presentation()

    # Use the 'Title Slide' layout (index 0)
    title_slide_layout = prs.slide_layouts[0]
    slide1 = prs.slides.add_slide(title_slide_layout)

    # Rename placeholders on the Title Slide
    slide1.shapes.title.name = 'Title'
    slide1.placeholders[1].name = 'Subtitle'

    # Use the 'Picture with Caption' layout (index 8) for content slides
    # This layout already has a title, a text box, and a picture placeholder.
    pic_caption_layout = prs.slide_layouts[8]

    # Add 6 content slides using this layout
    for _ in range(6):
        slide = prs.slides.add_slide(pic_caption_layout)

        # Rename the placeholders on each new slide for consistency
        # This is the key step.
        slide.shapes.title.name = 'Title'
        for shape in slide.placeholders:
            if 'Picture' in shape.name:
                shape.name = 'SlideImage'
            elif 'Text' in shape.name:
                shape.name = 'Content'

    prs.save(filepath)
    print(f"✅ Template '{filepath}' created successfully with named placeholders.")


if __name__ == "__main__":
    create_professional_template()