
"""
Auto Slide Deck Generator: Gemini + Web Search + Images -> PowerPoint

Dependencies:
  pip install google-generativeai python-pptx requests beautifulsoup4 pexels-api-python
"""
print("--- SCRIPT IS RUNNING ---")

#!/usr/bin/env python3

import argparse
import json
import os
import re
import sys
import time
import matplotlib.pyplot as plt
import io
from typing import List, Dict, Any
import requests
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
# -----------------------------
# Gemini Config
# -----------------------------
import google.generativeai as genai
genai.configure(api_key=os.getenv("GEMINI_API_KEY"))

# -----------------------------
# Pexels API Config
# -----------------------------
# CORRECTED IMPORT: The class is 'Pexels' and is imported directly.
# -----------------------------
# Pexels API Config
# -----------------------------
# -----------------------------
# Pexels API Config
# -----------------------------
# CORRECTED IMPORT: The error message confirmed the class is 'PexelsAPI'.
from pexels_api_python import PexelsAPI

try:
    # We are keeping the hardcoded key to avoid terminal issues.
    # PASTE YOUR KEY HERE if you haven't already.
    PEXELS_API_KEY = "NAZ3vZkeJ2yNZ9TwMrcoMi01D3Ysx7AieT7LHXa82MGc5wORqLKaoYOb"

    if not PEXELS_API_KEY or PEXELS_API_KEY == "YOUR_PEXELS_API_KEY_HERE":
        raise ValueError("Pexels API Key not provided in the script.")

    # CORRECTED USAGE: Use PexelsAPI() here to match the import statement.
    pexels_client = PexelsAPI(PEXELS_API_KEY)

except (ValueError, NameError) as e:
    print(f"[warn] {e}. Image generation will be skipped.", file=sys.stderr)
    pexels_client = None
# -----------------------------
# Constants
# -----------------------------
DEFAULT_RESULTS = 5
DEFAULT_TEMPLATE_NAME = "default_template.pptx"
USER_AGENT = (
    "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/123.0.0.0 Safari/537.36"
)
MAX_PER_PAGE_CHARS = 3000
REQUEST_TIMEOUT = 12


# -----------------------------
# Template & Image Functions
# -----------------------------
def create_default_template_if_not_exists(filepath: str):
    """Creates a simple template with text and image placeholders if one doesn't exist."""
    if os.path.exists(filepath):
        return
    print(f"Template not found. Creating a default template at: {filepath}")
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[0])  # Title
    prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content (Overview)
    prs.slides.add_slide(prs.slide_layouts[5])  # Title Only (for a Section Break)
    prs.slides.add_slide(prs.slide_layouts[3])  # Two Content (Key Point 1 & 2)
    prs.slides.add_slide(prs.slide_layouts[8])  # Picture with Caption (Key Point 3)
    prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content (Conclusion)
    prs.save(filepath)
    print("Default template with varied layouts created successfully.")


# MODIFIED: To be more flexible and detect all named shapes, not just placeholders.
# FINAL, CORRECTED version of analyze_template
# FINAL, UPGRADED version of analyze_template with table header extraction
def analyze_template(template_path: str) -> List[Dict[str, Any]]:
    """
    Analyzes a PowerPoint template and identifies all named shapes,
    including text, picture, and table placeholders, and extracts table headers.
    """
    prs = Presentation(template_path)
    slides_structure = []
    for i, slide in enumerate(prs.slides):
        placeholders = []
        for shape in slide.shapes:
            if not shape.name:
                continue

            placeholder_info = {"name": shape.name, "id": shape.shape_id}

            if hasattr(shape, 'table'):
                placeholder_info["type"] = "TABLE"
                # --- NEW: Extract header row from the template's table ---
                table = shape.table
                if table.rows:
                    placeholder_info["table_headers"] = [cell.text for cell in table.rows[0].cells]
                # --- End of New Logic ---
            elif shape.shape_type == 14:  # Picture placeholder type
                placeholder_info["type"] = "PICTURE"
            elif hasattr(shape, "text_frame"):
                placeholder_info["type"] = "TEXT"
            else:
                continue

            placeholders.append(placeholder_info)
        slides_structure.append({"slide_number": i + 1, "placeholders": placeholders})
    return slides_structure
# NEW: Function to search and download an image
# NEW: Function to search and download an image
# Function to search and download an image
def search_pexels_for_image(query: str) -> io.BytesIO | None:
    """Searches Pexels for a query and returns the image data as an in-memory stream."""
    if not pexels_client:
        return None
    try:
        search_results = pexels_client.search_photos(query=query, page=1, per_page=1)

        # CORRECTED RESULT ACCESS: The results are in a dictionary with a 'photos' key.
        photos = search_results.get('photos', [])

        if not photos:
            print(f"[warn] No Pexels image found for query: '{query}'", file=sys.stderr)
            return None

        # CORRECTED URL ACCESS: Each photo is a dictionary, and the URL is nested.
        photo_url = photos[0]['src']['large']

        response = requests.get(photo_url, timeout=REQUEST_TIMEOUT)
        response.raise_for_status()
        return io.BytesIO(response.content)

    except Exception as e:
        print(f"[error] Failed to fetch image from Pexels for query '{query}': {e}", file=sys.stderr)
        return None


# NEW: Function to create a chart image from data
def create_chart_image(chart_data: Dict[str, Any]) -> io.BytesIO | None:
    """Creates a bar or line chart from LLM-generated data using Matplotlib."""
    try:
        plt.style.use('seaborn-v0_8-talk')  # Use a professional style
        fig, ax = plt.subplots(figsize=(10, 5.5))

        chart_type = chart_data.get("type", "bar")
        labels = chart_data.get("labels", [])
        values = chart_data.get("values", [])

        if chart_type == 'bar':
            ax.bar(labels, values, color='#00A2FF')
        elif chart_type == 'line':
            ax.plot(labels, values, marker='o', linestyle='-', color='#0D47A1')

        ax.set_title(chart_data.get("title", ""), fontsize=16)
        ax.set_xlabel(chart_data.get("x_axis_label", ""), fontsize=12)
        ax.set_ylabel(chart_data.get("y_axis_label", ""), fontsize=12)
        ax.tick_params(axis='x', rotation=45)
        fig.tight_layout()

        # Save the plot to an in-memory buffer
        img_buffer = io.BytesIO()
        plt.savefig(img_buffer, format='png', dpi=150)
        plt.close(fig)
        img_buffer.seek(0)
        return img_buffer

    except Exception as e:
        print(f"[error] Failed to create chart: {e}", file=sys.stderr)
        return None
# -----------------------------
# Search and Scraping (Using SerpApi Only)
# -----------------------------
def web_search(query: str, count: int = DEFAULT_RESULTS) -> List[Dict[str, str]]:
    api_key = os.getenv("SERPAPI_API_KEY")
    if not api_key:
        raise ValueError("SERPAPI_API_KEY environment variable not set.")
    print("[info] Using SerpApi for web results.")
    url = "https://serpapi.com/search.json"
    params = {"engine": "google", "q": query, "api_key": api_key, "num": count}
    try:
        r = requests.get(url, params=params, timeout=REQUEST_TIMEOUT)
        r.raise_for_status()
        data = r.json()
        results = [{"title": it.get("title", "").strip(), "url": it.get("link", "").strip(),
                    "snippet": (it.get("snippet") or "").strip()} for it in data.get("organic_results", [])[:count]]
        if not results:
            raise RuntimeError("SerpApi returned no results.")
        return results
    except Exception as e:
        raise RuntimeError(f"SerpApi search failed: {e}. Please check your SERPAPI_API_KEY.")


def fetch_page_text(url: str) -> str:
    try:
        r = requests.get(url, headers={"User-Agent": USER_AGENT}, timeout=REQUEST_TIMEOUT)
        r.raise_for_status()
    except Exception:
        return ""  # Return empty string on failure
    soup = BeautifulSoup(r.text, "html.parser")
    for tag in soup(["script", "style", "nav", "footer", "aside"]):
        tag.decompose()
    return re.sub(r"\s+", " ", soup.get_text()).strip()[:MAX_PER_PAGE_CHARS]


# -----------------------------
# MODIFIED: LLM Integration (Gemini) with Image Descriptions
# -----------------------------
# MODIFIED: LLM Integration aware of custom layouts
# MODIFIED: LLM Integration with Speaker Notes
# FINAL: LLM Integration with Charts, Speaker Notes, and Dynamic Layouts
def llm_fill_slides(topic: str, template_structure: List[Dict[str, Any]], search_context: str) -> Dict[str, Any]:
    """
    Uses the LLM to fill a complex presentation structure, including text, image queries,
    speaker notes, and chart data for various dynamic layouts.
    """
    template_json = json.dumps(template_structure, indent=2)

    prompt = f"""
You are an expert presentation creator, data analyst, and speechwriter. Your task is to generate the complete content for a presentation on the topic: "{topic}".
You must use the provided web search results as your primary source of information to ensure the content is factual, current, and data-driven.

**Web Search Results (Context):**
---
{search_context}
---

**Presentation Structure:**
The presentation has the following structure, provided as a JSON object. You must fill it in.
---
{template_json}
---

**Your Task & Output Format:**
Return ONLY a valid JSON object that mirrors the input structure. For each slide, you must adhere to the following rules:

1.  **Speaker Notes**: For EVERY slide, you MUST add a `"speaker_notes"` field. This should be a short paragraph (2-4 sentences) with talking points, extra details, or an anecdote for the presenter. This text will NOT appear on the slide itself.

2.  **Placeholder Content**: For each item in the "placeholders" list:
    * **Text & Comparisons**: If the placeholder `type` is **"TEXT"** (e.g., for names like "Title", "Content", "ContentLeft", "ContentRight","Subtitle"), add a `"text"` field with the generated content.Make sure to write in bullet points(4-5 bullet points) in content placeholders only.For Title and SUbtitle give appropriate headings accordingly don't put images in their place.
        * **MODIFIED INSTRUCTION**: If you find `ContentLeft` and `ContentRight` on the same slide, you MUST generate comparative content. The content for each side must be a list of **up to five short bullet points**, and each bullet point must contain **no more than 5 words**.
    * **Timeline**: If you see placeholders named with the pattern `Timeline#_Title` and `Timeline#_Desc`, you MUST generate a title and description for a step in a process or a historical timeline. For example, for "The Scientific Method," the steps would be "Observation," "Hypothesis," "Experimentation."
    * **Images**: If the placeholder `type` is **"PICTURE"** and its name is NOT "ChartImage", add an `"image_query"` field with a concise, 2-5 word descriptive search query for a stock photo website.
    * **Charts**: If a placeholder is named **`ChartImage`**, you MUST attempt to generate data for a chart. Add a `"chart_data"` field containing a JSON object with the following keys: "type" (bar or line), "title", "x_axis_label", "y_axis_label", "labels" (a list of strings), and "values" (a list of numbers). 
   # ... inside the prompt string ...
     * **Tables**: If a placeholder `type` is **"TABLE"**, you MUST generate data for it. The required table headers are provided in the `"table_headers"` key. Your generated `"table_data"` **must have the exact same number of columns** and should follow the topics of the provided headers. The first row of your `"table_data"` must be the same as the provided `"table_headers"`.
**Example 1: A standard content slide**
{{
  "slide_number": 2,
  "placeholders": [
    {{
      "name": "Title", "type": "TEXT",
      "text": "Key Applications of AI"
    }},
    {{
      "name": "Content", "type": "TEXT",
      "text": "- Healthcare: Predictive diagnostics and personalized medicine.\\n- Finance: Algorithmic trading and fraud detection."
    }},
    {{
      "name": "SlideImage", "type": "PICTURE",
      "image_query": "artificial intelligence in medicine"
    }}
  ],
  "speaker_notes": "On this slide, we'll touch on two major areas where AI is making a significant impact. In healthcare, it's revolutionizing how we diagnose diseases. In finance, it's the backbone of modern trading systems."
}}

**Example 2: A chart slide**
{{
  "slide_number": 3,
  "placeholders": [
    {{
      "name": "Title", "type": "TEXT",
      "text": "Market Growth Projections"
    }},
    {{
      "name": "ChartImage", "type": "PICTURE",
      "chart_data": {{
        "type": "bar",
        "title": "Projected Growth of AI Market (in Billions USD)",
        "x_axis_label": "Year",
        "y_axis_label": "Market Size (USD)",
        "labels": ["2024", "2025", "2026", "2027"],
        "values": [250, 350, 500, 700]
      }}
    }}
  ],
  "speaker_notes": "As you can see from this chart, the market for AI is projected to grow exponentially over the next few years, reaching an estimated 700 billion dollars by 2027. This highlights the immense economic importance of the technology."
  
}}
**Example for a Table Slide:**
{{
  "slide_number": 5,
  "placeholders": [
    {{ "name": "Title", "type": "TEXT", "text": "Q3 Financial Metrics" }},
    {{ 
      "name": "DataTable", "type": "TABLE",
      "table_headers": ["Metric", "Q3 2023", "Q3 2024", "Change (%)"],
      "table_data": [
        ["Metric", "Q3 2023", "Q3 2024", "Change (%)"],
        ["Revenue", "$375.3M", "$242.2M", "-35.5%"],
        ["Net Loss", "($7.3M)", "($15.3M)", "-109.6%"]
      ]
    }}
  ],
  "speaker_notes": "This table provides a year-over-year comparison of our key financial metrics..."
}}
"""
    model = genai.GenerativeModel("gemini-1.5-flash")
    resp = model.generate_content(prompt)
    content = resp.text.strip()
    try:
        cleaned_json = re.sub(r"^```json\s*|\s*```$", "", content, flags=re.MULTILINE)
        return json.loads(cleaned_json)
    except (json.JSONDecodeError, AttributeError) as e:
        raise RuntimeError(f"Gemini returned non-JSON output. Error: {e}\nRaw Output:\n{content[:500]}")
# -----------------------------
# MODIFIED: PowerPoint Generation with Images
# -----------------------------
from pptx.dml.color import RGBColor
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx import Presentation

def style_text_frame(tf, is_title=False):
    """Apply consistent styling, wrapping, alignment, and safe font sizes."""
    tf.word_wrap = True
    tf.auto_size = False
    tf.margin_bottom = Pt(6)
    tf.margin_top = Pt(6)

    if is_title:
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE   # Titles centered
    else:
        tf.vertical_anchor = MSO_ANCHOR.TOP      # Body text starts at top

    for p in tf.paragraphs:
        p.alignment = PP_ALIGN.LEFT if not is_title else PP_ALIGN.CENTER
        for run in p.runs:
            run.font.name = "Calibri"
            run.font.color.rgb = RGBColor(40, 40, 40)
            if is_title:
                run.font.size = Pt(32)
                run.font.bold = True
            else:
                run.font.size = Pt(20)
                run.font.bold = False


def insert_image(slide, image_stream, shape_found):
    """Insert and resize image to fit within placeholder box."""
    pic = slide.shapes.add_picture(image_stream, shape_found.left, shape_found.top)

    # Max allowed size = placeholder box size (both width and height)
    max_width = shape_found.width
    max_height = shape_found.height

    # Scale proportionally to fit
    ratio = min(max_width / pic.width, max_height / pic.height)
    pic.width = int(pic.width * ratio)
    pic.height = int(pic.height * ratio)

    # Center inside the placeholder
    pic.left = shape_found.left + (shape_found.width - pic.width) // 2
    pic.top = shape_found.top + (shape_found.height - pic.height) // 2
    return pic

def set_cell_border(cell, border_color="444444", border_width='12700'):
    """Adds a solid border to all sides of a table cell."""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for border_name in ['a:lnL', 'a:lnR', 'a:lnT', 'a:lnB']: # Left, Right, Top, Bottom
        ln = OxmlElement(border_name)
        ln.set('w', border_width)
        ln.set('cap', 'flat')
        ln.set('cmpd', 'sng')
        ln.set('algn', 'ctr')
        solidFill = OxmlElement('a:solidFill')
        srgbClr = OxmlElement('a:srgbClr')
        srgbClr.set('val', border_color)
        solidFill.append(srgbClr)
        ln.append(solidFill)
        tcPr.append(ln)
# MODIFIED: PowerPoint Generation with Images and Speaker Notes
# MODIFIED: PowerPoint Generation with Charts, Images, and Speaker Notes
def populate_ppt(template_path: str, filled_slides: List[Dict[str, Any]], outfile: str,
                 references: List[Dict[str, str]]):
    prs = Presentation(template_path)
    for slide_data in filled_slides:
        slide_idx = slide_data["slide_number"] - 1
        if slide_idx >= len(prs.slides):
            continue
        slide = prs.slides[slide_idx]

        speaker_notes_text = slide_data.get("speaker_notes")
        if speaker_notes_text:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = speaker_notes_text

        placeholders_to_remove = []
        for placeholder_data in slide_data.get("placeholders", []):
            placeholder_name = placeholder_data.get("name")
            shape_found = next((shape for shape in slide.shapes if shape.name == placeholder_name), None)

            if not shape_found: continue

            if placeholder_data.get("text"):
                shape_found.text_frame.clear()
                shape_found.text_frame.text = placeholder_data["text"]
                shape_found.text_frame.auto_size = True

            # --- NEW: Chart Generation Logic ---
            elif placeholder_data.get("chart_data"):
                print(f"  - Generating chart: '{placeholder_data['chart_data'].get('title', 'Untitled')}'")
                chart_image_stream = create_chart_image(placeholder_data["chart_data"])
                if chart_image_stream:
                    slide.shapes.add_picture(chart_image_stream, shape_found.left, shape_found.top,
                                             width=shape_found.width, height=shape_found.height)
                    placeholders_to_remove.append(shape_found)
            # --- End of New Logic ---
                    # --- NEW: Table Population Logic with Professional Styling ---
                    # --- FINAL: Table Population Logic with Grid Lines and Styling ---
            elif table_data := placeholder_data.get("table_data"):
                if not hasattr(shape_found, "table"): continue
                print(f"  - Populating and styling table...")

                table = shape_found.table

                # Set column widths for better spacing
                table.columns[0].width = Inches(3.5)  # Metric column
                table.columns[1].width = Inches(2.0)  # Change %
                table.columns[2].width = Inches(2.5)  # Q3 2024
                table.columns[3].width = Inches(2.5)  # Q3 2023

                # Populate data and apply styles cell by cell
                for i, row_data in enumerate(table_data):
                    for j, cell_text in enumerate(row_data):
                        if i < len(table.rows) and j < len(table.columns):
                            cell = table.cell(i, j)
                            cell.text = str(cell_text)
                            cell.vertical_anchor = MSO_ANCHOR.MIDDLE  # Center text vertically

                            # Add a border to every cell
                            set_cell_border(cell)

                            p = cell.text_frame.paragraphs[0]
                            if not p.runs: p.add_run()
                            run = p.runs[0]
                            run.font.name = 'Calibri'
                            run.font.size = Pt(14)

                            # Style Header Row
                            if i == 0:
                                p.alignment = PP_ALIGN.CENTER
                                run.font.bold = True
                                cell.fill.solid()
                                cell.fill.fore_color.rgb = RGBColor(0x44, 0x54, 0x6A)  # Dark blue-grey
                                run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)  # White text
                            # Style Data Rows
                            else:
                                # Right-align all columns except the first one (Metric)
                                p.alignment = PP_ALIGN.RIGHT if j > 0 else PP_ALIGN.LEFT
                                run.font.bold = False
            elif placeholder_data.get("image_query"):
                print(f"  - Searching for image: '{placeholder_data['image_query']}'")
                image_stream = search_pexels_for_image(placeholder_data["image_query"])
                if image_stream:
                    slide.shapes.add_picture(image_stream, shape_found.left, shape_found.top, width=shape_found.width,
                                             height=shape_found.height)
                    placeholders_to_remove.append(shape_found)

        for shape in placeholders_to_remove:
            sp = shape.element
            sp.getparent().remove(sp)

    # Add References slide
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "References"
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for ref in references:
        p = tf.add_paragraph()
        p.text = f"{ref['title']}"
        p.font.size = Pt(12)
        run = p.add_run()
        run.text = f" ({ref['url']})"
        run.font.size = Pt(9)
        run.hyperlink.address = ref['url']
    prs.save(outfile)
# Function to classify the topic into a category
def classify_topic(topic: str) -> str:
    """Uses the LLM to classify a topic into a predefined category."""
    print("[info] Classifying topic...")

    # These categories MUST match your template filenames (without .pptx)
    categories = ["technology", "business", "science"]

    prompt = f"""
    Classify the following topic into one of these categories: {', '.join(categories)}.
    Topic: "{topic}"
    Respond with only a single word from the list. If it doesn't fit, respond with 'business'.
    """

    try:
        model = genai.GenerativeModel("gemini-1.5-flash")
        response = model.generate_content(prompt)
        category = response.text.strip().lower()

        if category in categories:
            print(f"[info] Topic classified as: {category}")
            return category
        else:
            print(f"[warn] Classification unclear. Defaulting to 'business'.")
            return "business"
    except Exception as e:
        print(f"[error] LLM classification call failed: {e}. Defaulting to 'business'.")
        return "business"

# -----------------------------
# Main
# -----------------------------
# MODIFIED: Main function simplified for single-template testing
# FINAL Main function with intelligent template selection
def main():
    parser = argparse.ArgumentParser(
        description="Generate a slide deck with images from a topic using a template, Gemini, and Web Search.")
    parser.add_argument("--topic", type=str, required=True, help="The topic for the presentation.")

    # Use --template_dir to point to your folder of templates
    parser.add_argument("--template_dir", type=str, default="templates",
                        help="Path to the directory containing your .pptx template files.")

    parser.add_argument("--results", type=int, default=DEFAULT_RESULTS, help="Number of web search results to use.")
    parser.add_argument("--outfile", type=str, default=None, help="Name of the output file.")
    args = parser.parse_args()

    topic = args.topic.strip()
    template_dir = args.template_dir
    outname = args.outfile or (re.sub(r"[^A-Za-z0-9_-]+", "_", topic) + ".pptx")

    # --- Intelligent Template Selection Logic ---
    topic_category = classify_topic(topic)
    template_path = os.path.join(template_dir, f"{topic_category}.pptx")

    # Fallback to the business template if the specific one doesn't exist
    if not os.path.exists(template_path):
        print(f"[warn] Template '{template_path}' not found. Falling back to the business template.")
        template_path = os.path.join(template_dir, "business.pptx")
        if not os.path.exists(template_path):
            raise FileNotFoundError(
                f"Default template 'business.pptx' not found in '{template_dir}'. Please create it.")

    # The rest of the process remains the same
    print(f"[1/6] Using template: {template_path}")
    search_results = web_search(topic, count=args.results)

    print(f"[2/6] Fetching text from top results...")
    context_parts = [f"Source: {item['title']}\nURL: {item['url']}\nContent: {fetch_page_text(item['url'])}" for item in
                     search_results]

    print(f"[3/6] Analyzing PowerPoint template...")
    template_structure = analyze_template(template_path)

    print("[4/6] Calling Gemini to generate slide content and image queries...")
    filled_slides_data = llm_fill_slides(topic, template_structure, "\n---\n".join(context_parts))

    print(f"[5/6] Populating PowerPoint with text and images...")
    populate_ppt(template_path, filled_slides_data, outname, search_results)

    print(f"[6/6] Done! ✅ Presentation saved to: {outname}")
if __name__ == "__main__":
    main()