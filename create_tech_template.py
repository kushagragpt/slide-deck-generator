# create_final_template.py
from pptx import Presentation
from pptx.util import Inches

def create_final_template(filepath="final_template.pptx"):
    """
    Generates a final, robust 7-slide template using the most reliable
    built-in layouts and correctly named placeholders.
    """
    print(f"Creating a final, robust template at: {filepath}")
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title Slide (Layout 0)
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.name = 'Title'
    slide.placeholders[1].name = 'Subtitle'

    # Slide 2: Overview (Layout 8 - Picture w/ Caption)
    slide = prs.slides.add_slide(prs.slide_layouts[8])
    slide.shapes.title.name = 'Title'
    for shape in slide.placeholders:
        if 'Picture' in shape.name: shape.name = 'SlideImage'
        elif 'Text' in shape.name: shape.name = 'Content'

    # Slide 3: Key Point 1 (Layout 8 - Picture w/ Caption)
    slide = prs.slides.add_slide(prs.slide_layouts[8])
    slide.shapes.title.name = 'Title'
    for shape in slide.placeholders:
        if 'Picture' in shape.name: shape.name = 'SlideImage'
        elif 'Text' in shape.name: shape.name = 'Content'

    # Slide 4: Comparison Slide (Layout 3 - Comparison)
    slide = prs.slides.add_slide(prs.slide_layouts[3])
    slide.shapes.title.name = 'Title'
    # This layout has two content placeholders by default
    slide.placeholders[1].name = 'ContentLeft'
    slide.placeholders[2].name = 'ContentRight'

    # Slide 5: Chart Slide (Layout 1 - Title and Content)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.name = 'Title'
    # We'll use the main content area for the chart
    slide.placeholders[1].name = 'ChartImage'

    # Slide 6: Key Point 4 (Layout 8 - Picture w/ Caption)
    slide = prs.slides.add_slide(prs.slide_layouts[8])
    slide.shapes.title.name = 'Title'
    for shape in slide.placeholders:
        if 'Picture' in shape.name: shape.name = 'SlideImage'
        elif 'Text' in shape.name: shape.name = 'Content'

    # Slide 7: Conclusion (Layout 8 - Picture w/ Caption)
    slide = prs.slides.add_slide(prs.slide_layouts[8])
    slide.shapes.title.name = 'Title'
    for shape in slide.placeholders:
        if 'Picture' in shape.name: shape.name = 'SlideImage'
        elif 'Text' in shape.name: shape.name = 'Content'

    prs.save(filepath)
    print(f"✅ Template '{filepath}' created successfully.")

if __name__ == "__main__":
    create_final_template()