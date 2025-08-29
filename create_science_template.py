# create_science_template.py
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


def create_science_template(filepath="science.pptx"):
    """
    Generates a final, robust 7-slide science-themed PowerPoint template that includes
    layouts for content, timelines, and charts.
    """
    print(f"Creating a new 7-slide science-themed template at: {filepath}")
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Define Color Palette (Science Theme) ---
    BACKGROUND_COLOR = RGBColor(0xFF, 0xFF, 0xFF)  # White
    ACCENT_COLOR = RGBColor(0x00, 0x69, 0x5C)  # Scientific Teal
    BODY_TEXT_COLOR = RGBColor(0x2F, 0x2F, 0x2F)  # Very Dark Grey

    # --- Function to apply background color ---
    def apply_background(slide):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = BACKGROUND_COLOR

    # --- Slide 1: Title Slide (Robust Method) ---
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    apply_background(slide)
    # This is the most reliable way to name the title/subtitle placeholders
    slide.shapes.title.name = 'Title'
    slide.placeholders[1].name = 'Subtitle'
    # Optional styling
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = ACCENT_COLOR
    slide.placeholders[1].text_frame.paragraphs[0].font.color.rgb = BODY_TEXT_COLOR

    # --- Slide 2: Standard Content (Abstract / Overview) ---
    slide = prs.slides.add_slide(prs.slide_layouts[8])  # Picture with Caption
    apply_background(slide)
    slide.shapes.title.name = 'Title'
    for shape in slide.placeholders:
        if 'Picture' in shape.name:
            shape.name = 'SlideImage'
        elif 'Text' in shape.name:
            shape.name = 'Content'

    # --- Slide 3: Timeline / Process Slide ---
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
    apply_background(slide)
    slide.shapes.title.name = 'Title'
    # Add named textboxes for a 3-step process
    step_width = Inches(4.0)
    step_margin = Inches(0.41)
    slide.shapes.add_textbox(Inches(0.5), Inches(1.5), step_width, Inches(1)).name = 'Timeline1_Title'
    slide.shapes.add_textbox(Inches(0.5), Inches(2.5), step_width, Inches(4)).name = 'Timeline1_Desc'
    slide.shapes.add_textbox(Inches(0.5) + step_width + step_margin, Inches(1.5), step_width,
                             Inches(1)).name = 'Timeline2_Title'
    slide.shapes.add_textbox(Inches(0.5) + step_width + step_margin, Inches(2.5), step_width,
                             Inches(4)).name = 'Timeline2_Desc'
    slide.shapes.add_textbox(Inches(0.5) + (step_width + step_margin) * 2, Inches(1.5), step_width,
                             Inches(1)).name = 'Timeline3_Title'
    slide.shapes.add_textbox(Inches(0.5) + (step_width + step_margin) * 2, Inches(2.5), step_width,
                             Inches(4)).name = 'Timeline3_Desc'

    # --- Slide 4: Chart Slide (Results) ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    apply_background(slide)
    slide.shapes.title.name = 'Title'
    slide.placeholders[1].name = 'ChartImage'  # Rename the content placeholder for the chart

    # --- Slides 5, 6, 7: Standard Content ---
    for _ in range(3):  # Discussion, Conclusion, etc.
        slide = prs.slides.add_slide(prs.slide_layouts[8])  # Picture with Caption
        apply_background(slide)
        slide.shapes.title.name = 'Title'
        for shape in slide.placeholders:
            if 'Picture' in shape.name:
                shape.name = 'SlideImage'
            elif 'Text' in shape.name:
                shape.name = 'Content'

    prs.save(filepath)
    print(f"✅ Template '{filepath}' with dynamic layouts created successfully.")


if __name__ == "__main__":
    create_science_template()